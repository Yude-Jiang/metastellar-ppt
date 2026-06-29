"""ST brand helper library for python-pptx.

Encodes the rules from skills/st-ppt-brand (palette, typography, message bar,
Title-Only content slides, card rows, and diagram primitives) so generated decks
are on-brand by construction. Import this from the deck-build script.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- Official ST palette (Green Vogue / Gold / Picton Blue / White) ----
ST_DARK_BLUE  = RGBColor(0x03, 0x23, 0x4B)   # Green Vogue
ST_YELLOW     = RGBColor(0xFF, 0xD2, 0x00)   # Gold
ST_LIGHT_BLUE = RGBColor(0x3C, 0xB4, 0xE6)   # Picton Blue
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_1        = RGBColor(0xEE, 0xEF, 0xF1)
GRAY_2        = RGBColor(0xDB, 0xDE, 0xE1)
GRAY_3        = RGBColor(0xC0, 0xC8, 0xD2)
SLATE         = RGBColor(0x42, 0x59, 0x78)   # first shade of dark blue

# Dark-blue shading ramp (graded headers / process steps), darkest first
RAMP = [ST_DARK_BLUE, SLATE, RGBColor(0x80, 0x91, 0xA5), RGBColor(0xC0, 0xC9, 0xCE)]

FONT = "Arial"
SLIDE_W = 13.333
SLIDE_H = 7.5


# Dark fills: white text. Light fills: ST Dark Blue text.
# RAMP[2] (#8091A5) is mid-tone — white text for readable contrast.
_DARK_FILLS = frozenset({ST_DARK_BLUE, SLATE, RAMP[2]})


def ramp_text(step):
    return text_on(RAMP[min(step, 3)])


def text_on(fill_color):
    """Mandatory contrast: white on dark fills; ST Dark Blue on light fills."""
    if fill_color in _DARK_FILLS:
        return WHITE
    return ST_DARK_BLUE


BODY_SIZE = 14   # prefer 14 pt per brand spec (12 min, 20 max)
TITLE_SIZE = 27
MSG_BAR_SIZE = 20


def new_deck():
    """16:9 presentation per ST template geometry."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def title_only_slide(prs):
    """Reproduce the 'Title Only' layout: blank slide; place shapes yourself."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _place_logo(slide, logo_path, left, top, height=0.48):
    if logo_path:
        slide.shapes.add_picture(logo_path, Inches(left), Inches(top), height=Inches(height))


def presentation_title_slide(prs, title, presenter=None, logo_path=None):
    """Main / presentation title — navy field, left yellow accent, white title (see special-slides.md)."""
    slide = title_only_slide(prs)
    _slide_bg(slide, ST_DARK_BLUE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    Inches(0.14), Inches(SLIDE_H))
    fill(accent, ST_YELLOW)
    tb = slide.shapes.add_textbox(Inches(1.1), Inches(2.55), Inches(9.5), Inches(1.15))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    _style(tf, WHITE, 36, bold=True, align=PP_ALIGN.LEFT)
    if presenter:
        ptb = slide.shapes.add_textbox(Inches(1.1), Inches(3.85), Inches(8.0), Inches(0.55))
        ptf = ptb.text_frame
        no_autofit(ptf)
        ptf.text = presenter
        _style(ptf, WHITE, 18, bold=False, align=PP_ALIGN.LEFT)
    _place_logo(slide, logo_path, 11.35, 0.38, height=0.5)
    return slide


def agenda_slide(prs, topics, title="Agenda", logo_path=None, columns=2):
    """Agenda / table of contents — white field, yellow numbered tiles (see special-slides.md)."""
    slide = title_only_slide(prs)
    _slide_bg(slide, WHITE)
    tb = slide.shapes.add_textbox(Inches(9.6), Inches(0.32), Inches(3.4), Inches(0.9))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _style(tf, ST_DARK_BLUE, 32, bold=True, align=PP_ALIGN.RIGHT)
    n = len(topics)
    cols = max(1, min(columns, 2))
    rows = (n + cols - 1) // cols
    col_x = [1.35, 7.05]
    row_y0 = 1.55
    row_gap = 1.05
    tile = 0.42
    for i, topic in enumerate(topics):
        col = i // rows if rows else 0
        row = i % rows if rows else 0
        if col >= cols:
            col = cols - 1
        x = col_x[col]
        y = row_y0 + row * row_gap
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                    Inches(tile), Inches(tile))
        fill(sq, ST_YELLOW)
        stf = sq.text_frame
        no_autofit(stf)
        stf.vertical_anchor = MSO_ANCHOR.MIDDLE
        stf.text = str(i + 1)
        for p in stf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        _style(stf, ST_DARK_BLUE, 14, bold=True, align=PP_ALIGN.CENTER)
        label = slide.shapes.add_textbox(Inches(x + tile + 0.18), Inches(y + 0.06),
                                         Inches(4.8), Inches(0.38))
        ltf = label.text_frame
        no_autofit(ltf)
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ltf.text = topic
        _style(ltf, ST_DARK_BLUE, 14, bold=False, align=PP_ALIGN.LEFT)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.42)
    return slide


def section_title_slide(prs, title, image_path=None, logo_path=None):
    """Section break — navy field, top yellow bar, optional hero image (see special-slides.md)."""
    slide = title_only_slide(prs)
    _slide_bg(slide, ST_DARK_BLUE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.35), Inches(0),
                                 Inches(SLIDE_W - 3.35), Inches(1.15))
    fill(bar, ST_YELLOW)
    tf = bar.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.35)
    tf.text = title
    _style(tf, ST_DARK_BLUE, 30, bold=True, align=PP_ALIGN.LEFT)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(2.2), Inches(1.65),
                                 width=Inches(8.9), height=Inches(4.35))
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(6.72),
                                  Inches(12.4), Inches(0.02))
    fill(rule, GRAY_3)
    _place_logo(slide, logo_path, 0.45, 6.85, height=0.38)
    return slide


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def no_autofit(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def _style(tf, color, size_pt, bold=False, font=FONT, align=None):
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size_pt)
            r.font.bold = bold
            r.font.color.rgb = color


def corner_accent(slide):
    """Thin ST Dark Blue block, top-right (template accent)."""
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.55), Inches(0.0),
                               Inches(0.78), Inches(0.34))
    fill(a, ST_DARK_BLUE)
    return a


def add_title(slide, text, subtitle=None, align=PP_ALIGN.RIGHT, size=27):
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(0.28), Inches(12.1), Inches(0.95))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = text
    tf.paragraphs[0].alignment = align
    _style(tf, ST_DARK_BLUE, size, bold=True)
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.alignment = align
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(13)
            r.font.color.rgb = ST_LIGHT_BLUE
    return tb


def add_message_bar(slide, text, fill_color=ST_LIGHT_BLUE):
    """The signature ST element. Geometry fixed by the template; 20pt Arial only.
    Fill must be ST Yellow / Dark Blue / Light Blue / slate. Text is single color.
    """
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.43),
                                 Inches(9.8), Inches(0.84))
    fill(bar, fill_color)
    tf = bar.text_frame
    no_autofit(tf)
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = text
    txt = text_on(fill_color)
    _style(tf, txt, MSG_BAR_SIZE, bold=True)
    return bar


def box(slide, x, y, w, h, text, fill_color, text_color=None, size=BODY_SIZE, bold=True,
        align=PP_ALIGN.CENTER, sub=None, sub_size=12):
    """A rectangle with embedded text. text_color defaults to text_on(fill)."""
    if text_color is None:
        text_color = text_on(fill_color)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    fill(sh, fill_color)
    tf = sh.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.text = text
    p0 = tf.paragraphs[0]
    p0.alignment = align
    for r in p0.runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = text_color
    if sub:
        p = tf.add_paragraph()
        p.text = sub
        p.alignment = align
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(sub_size)
            r.font.color.rgb = text_color
    return sh


def bullet_box(slide, x, y, w, h, bullets, shade=GRAY_1, size=BODY_SIZE, heading=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    fill(sh, shade)
    tf = sh.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_right = Inches(0.14)
    first = True
    if heading:
        tf.text = heading
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        for r in tf.paragraphs[0].runs:
            r.font.name = FONT
            r.font.size = Pt(size + 2)
            r.font.bold = True
            r.font.color.rgb = ST_DARK_BLUE
        first = False
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "\u2022 " + b
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = ST_DARK_BLUE
    return sh


def label(slide, x, y, w, text, color=ST_DARK_BLUE, size=11, bold=True,
          align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    tf = tb.text_frame
    no_autofit(tf)
    tf.margin_left = Inches(0.02)
    tf.margin_top = Inches(0)
    tf.text = text
    tf.paragraphs[0].alignment = align
    for r in tf.paragraphs[0].runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def arrow(slide, x1, y1, x2, y2, color=ST_DARK_BLUE, width=2.25, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),
                             {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if dashed:
        ln.insert(0, ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return c


def dashed_container(slide, x, y, w, h, color=SLATE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = Pt(1.5)
    ln = sh.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    sh.shadow.inherit = False
    return sh


def add_activation_timeline(
    slide,
    checkpoints,
    top_items=None,
    bottom_items=None,
    x0=0.55,
    x1=12.8,
    y_line=3.95,
):
    """Draw a 2-lane activation timeline similar to GTM launch plans.

    checkpoints: [{"x": 2.4, "label": "Mar 5"}, ...]
    top_items / bottom_items:
      [{"x": 2.6, "title": "Milestone", "note": "optional", "w": 2.0, "color": ST_YELLOW}, ...]
    """
    top_items = top_items or []
    bottom_items = bottom_items or []

    # Main axis
    arrow(slide, x0, y_line, x1, y_line, color=ST_DARK_BLUE, width=2.2)

    # Lane labels
    box(slide, 0.0, y_line - 1.55, 0.42, 1.42, "CONTENT", ST_DARK_BLUE, WHITE, size=10)
    box(slide, 0.0, y_line + 0.15, 0.42, 1.42, "PROMO", ST_DARK_BLUE, WHITE, size=10)

    # Checkpoints on axis
    for cp in checkpoints:
        x = cp["x"]
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.08), Inches(y_line - 0.08),
                                     Inches(0.16), Inches(0.16))
        fill(dot, WHITE)
        dot.line.color.rgb = ST_DARK_BLUE
        dot.line.width = Pt(1.5)
        if cp.get("label"):
            label(slide, x - 0.35, y_line - 0.42, 0.8, cp["label"], color=ST_DARK_BLUE,
                  size=10, bold=False, align=PP_ALIGN.CENTER)

    # Activity cards above axis
    for i, item in enumerate(top_items):
        w = item.get("w", 2.0)
        h = 0.5 if not item.get("note") else 0.82
        y = y_line - 1.2 - (0.62 * (i % 2))
        fc = item.get("color", ST_YELLOW)
        tc = text_on(fc)
        box(slide, item["x"] - w / 2, y, w, h, item["title"], fc, tc, size=11,
            align=PP_ALIGN.LEFT, sub=item.get("note"), sub_size=10)
        arrow(slide, item["x"], y + h, item["x"], y_line - 0.12, color=GRAY_3, width=1.2)

    # Activity cards below axis
    for i, item in enumerate(bottom_items):
        w = item.get("w", 2.0)
        h = 0.5 if not item.get("note") else 0.82
        y = y_line + 0.26 + (0.62 * (i % 2))
        fc = item.get("color", GRAY_1)
        tc = text_on(fc)
        box(slide, item["x"] - w / 2, y, w, h, item["title"], fc, tc, size=11,
            align=PP_ALIGN.LEFT, sub=item.get("note"), sub_size=10)
        arrow(slide, item["x"], y, item["x"], y_line + 0.12, color=GRAY_3, width=1.2)


def timeline_template_slide(prs, title, message, checkpoints, top_items=None, bottom_items=None):
    """One-call template: title + message bar + 2-lane timeline."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    add_title(slide, title, align=PP_ALIGN.RIGHT, size=TITLE_SIZE + 2)
    add_message_bar(slide, message, fill_color=ST_LIGHT_BLUE)
    add_activation_timeline(slide, checkpoints, top_items=top_items, bottom_items=bottom_items)
    return slide


def add_cards_row(slide, cards, top=2.5, bottom=6.95, gap=0.4, left_margin=0.45,
                  header="yellow", img_ratio=2.35):
    """cards = [{"title", "bullets":[...], "img": path|None}].
    header = 'yellow' (yellow bar, dark-blue text) or 'ramp' (graded dark-blue).
    Lays out N equal cards: header bar + optional image banner + gray bullet box.
    """
    from PIL import Image
    n = len(cards)
    total_w = SLIDE_W - 2 * left_margin
    w = (total_w - gap * (n - 1)) / n
    head_h = 0.5
    img_top = top + head_h + 0.08
    img_h = w / img_ratio
    box_top = img_top + img_h + 0.08
    for i, c in enumerate(cards):
        x = left_margin + i * (w + gap)
        if header == "yellow":
            hf, ht = ST_YELLOW, ST_DARK_BLUE
        else:
            hf, ht = RAMP[min(i, 3)], text_on(RAMP[min(i, 3)])
        box(slide, x, top, w, head_h, c["title"], hf, ht, size=BODY_SIZE,
            align=PP_ALIGN.LEFT)
        bt = box_top
        img = c.get("img")
        if img:
            # crop to a clean banner, then place with locked aspect ratio
            im = Image.open(img)
            iw, ih = im.size
            crop_h = min(ih, int(iw / img_ratio))
            cropped = img + ".banner.png"
            im.crop((0, 0, iw, crop_h)).save(cropped)
            slide.shapes.add_picture(cropped, Inches(x), Inches(img_top),
                                     width=Inches(w), height=Inches(img_h))
        else:
            bt = top + head_h
        bullet_box(slide, x, bt, w, top + (bottom - top) - bt,
                   c.get("bullets", []))


def footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.32))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = text
    for r in tf.paragraphs[0].runs:
        r.font.name = FONT
        r.font.size = Pt(9)
        r.font.color.rgb = ST_DARK_BLUE
    return tb


CLOSING_FOOTER = (
    "© STMicroelectronics - All rights reserved.\n"
    "ST logo is a trademark or a registered trademark of STMicroelectronics "
    "International NV or its affiliates in the EU and/or other countries.\n"
    "For additional information about ST trademarks, please refer to "
    "www.st.com/trademarks.\n"
    "All other product or service names are the property of their respective owners."
)


def closing_slide(prs, logo_path=None, tagline="Our technology starts with You"):
    """Mandatory ST closing / trademark slide for external decks (see compliance checklist)."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(5.85)
    )
    fill(panel, ST_DARK_BLUE)
    if logo_path:
        slide.shapes.add_picture(logo_path, Inches(0.55), Inches(0.45), height=Inches(0.55))
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(11.7), Inches(1.4))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = tagline
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _style(tf, WHITE, 32, bold=True)
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.95), Inches(SLIDE_W), Inches(1.55)
    )
    fill(band, ST_YELLOW)
    ftb = slide.shapes.add_textbox(Inches(0.45), Inches(6.05), Inches(12.4), Inches(1.35))
    ftf = ftb.text_frame
    no_autofit(ftf)
    ftf.word_wrap = True
    ftf.text = CLOSING_FOOTER
    for p in ftf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(8)
            r.font.color.rgb = ST_DARK_BLUE
    return slide
