"""Metastellar slide theme library for python-pptx.

Digital Enterprise Deck System — dashboard-driven enterprise aesthetic.
Blue mono palette, Arial typography, § chrome, KPI cards, MetaStellar mark.
Import from the deck-build script; rules live in skills/metastellar-slides/.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- Design tokens (Digital Enterprise Deck System) ----
# Ink
INK           = RGBColor(0x0B, 0x12, 0x20)
INK_2         = RGBColor(0x1E, 0x29, 0x3B)
INK_3         = RGBColor(0x47, 0x55, 0x69)
INK_4         = RGBColor(0x94, 0xA3, 0xB8)
INK_5         = RGBColor(0xCB, 0xD5, 0xE1)
# Paper
PAPER         = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_2       = RGBColor(0xF8, 0xFA, 0xFC)
PAPER_3       = RGBColor(0xF1, 0xF5, 0xF9)
LINE          = RGBColor(0xE2, 0xE8, 0xF0)
LINE_2        = RGBColor(0xCB, 0xD5, 0xE1)
# Brand (blue mono — single hue only)
BRAND         = RGBColor(0x1E, 0x40, 0xAF)
BRAND_2       = RGBColor(0x25, 0x63, 0xEB)
BRAND_3       = RGBColor(0x3B, 0x82, 0xF6)
BRAND_4       = RGBColor(0x60, 0xA5, 0xFA)
BRAND_DEEP    = RGBColor(0x0F, 0x1E, 0x4A)   # dark surfaces — never pure black
BRAND_PALE    = RGBColor(0xDB, 0xEA, 0xFE)
BRAND_TINT    = RGBColor(0xEF, 0xF6, 0xFF)
BRAND_GLOW    = RGBColor(0x93, 0xC5, 0xFD)
OK            = RGBColor(0x10, 0xB9, 0x81)   # live dots only
WHITE         = PAPER

# Back-compat aliases (legacy builder names)
BLUE_500      = BRAND_3
BLUE_600      = BRAND_2
BLUE_50       = BRAND_TINT
BLUE_100      = BRAND_PALE
BLUE_800      = BRAND
PRIMARY_DARK  = BRAND
ACCENT        = BRAND_3
ACCENT_LIGHT  = BRAND_3
GRAY_50       = PAPER_2
GRAY_100      = PAPER_3
GRAY_200      = LINE
GRAY_300      = LINE_2
GRAY_600      = INK_3
GRAY_800      = INK_2
GRAY_1        = PAPER_3
GRAY_2        = LINE
GRAY_3        = LINE_2
SLATE         = BRAND_2
RAMP = [BRAND_DEEP, BRAND, BRAND_2, BRAND_3]

FONT          = "Arial"
FONT_CN       = "Microsoft YaHei"
FONT_MONO     = "Consolas"
SLIDE_W       = 13.333
SLIDE_H       = 7.5
CHROME_PAD    = 0.67   # ~96px at 1920 — slide chrome inset
DEFAULT_TOTAL = 24     # SEO/GEO deck page count

# MetaStellar SEO/GEO defaults
SEO_GEO_TAGLINE = "The new science of visibility."
SEO_GEO_MANIFESTO = "If you can't be found, you don't exist."
SEO_GEO_CONTACT = "Let's make you findable."
DEFAULT_SEO_METRICS = [
    {"label": "Projects", "value": "32", "badge": "↑ 22% QoQ"},
    {"label": "Brands", "value": "120", "desc": "since 2020"},
    {"label": "Organic lift", "value": "4.2", "unit": "×", "desc": "12-mo avg"},
    {"label": "AI citation", "value": "68", "unit": "%", "desc": "4 engines"},
]

# Dark fills: white text. Light fills: BRAND text.
_DARK_FILLS = frozenset({BRAND_DEEP, BRAND, BRAND_2, INK, INK_2})


def ramp_text(step):
    return text_on(RAMP[min(step, 3)])


def text_on(fill_color):
    """Mandatory contrast: white on dark fills; BRAND on light fills."""
    if fill_color in _DARK_FILLS:
        return WHITE
    return BRAND


BODY_SIZE = 14
TITLE_SIZE = 27
MSG_BAR_SIZE = 20
KPI_VALUE_SIZE = 44
MONO_LABEL_SIZE = 10


def _style(tf, color, size_pt, bold=False, font=FONT, align=None, mono=False):
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.name = FONT_MONO if mono else font
            r.font.size = Pt(size_pt)
            r.font.bold = bold
            r.font.color.rgb = color
            if mono:
                r.font.name = FONT_MONO


def draw_metastellar_mark(slide, cx, cy, size_in, color=BRAND):
    """MetaStellar symbol — twin orbits + open triangle (3 strokes)."""
    s = size_in
    half = s / 2
    for rot in (-28, 28):
        ov = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - half), Inches(cy - half * 0.35),
            Inches(s), Inches(s * 0.35),
        )
        ov.rotation = rot
        ov.fill.background()
        ov.line.color.rgb = color
        ov.line.width = Pt(2)
    # Triangle strokes (open peak)
    lx, rx = cx - half * 0.38, cx + half * 0.38
    top_y, bot_y = cy - half * 0.36, cy + half * 0.32
    mid_y = cy + half * 0.18
    for x1, y1, x2, y2 in (
        (lx, bot_y, cx - half * 0.06, top_y),
        (cx + half * 0.06, top_y, rx, bot_y),
        (lx + half * 0.12, mid_y, rx - half * 0.12, mid_y),
    ):
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        ln.line.color.rgb = color
        ln.line.width = Pt(2)


def add_wordmark(slide, left, top, dark=False, compact=False):
    """MetaStellar lockup — mark + MetaStellar / 智元星启."""
    mark_h = 0.32 if compact else 0.42
    draw_metastellar_mark(slide, left + mark_h / 2, top + mark_h / 2, mark_h, WHITE if dark else BRAND)
    tx = left + mark_h + 0.14
    en_size = 11 if compact else 14
    tb = slide.shapes.add_textbox(Inches(tx), Inches(top), Inches(2.8), Inches(0.55 if not compact else 0.4))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = "MetaStellar"
    _style(tf, WHITE if dark else INK, en_size, bold=True, align=PP_ALIGN.LEFT)
    if not compact:
        p = tf.add_paragraph()
        p.text = "智元星启"
        for r in p.runs:
            r.font.name = FONT_CN
            r.font.size = Pt(9)
            r.font.color.rgb = WHITE if dark else INK_3


def add_slide_header(slide, section_num, section_name, section_title=None,
                     section_cn=None, rhs_meta=None, dark=False):
    """§ chapter pill + section title + optional rhs meta (every inner slide)."""
    y = CHROME_PAD * 0.45
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CHROME_PAD), Inches(y), Inches(1.55), Inches(0.28),
    )
    fill(pill, BRAND_TINT if not dark else BRAND)
    pill.line.color.rgb = BRAND_PALE if not dark else BRAND_2
    pill.line.width = Pt(1)
    ptf = pill.text_frame
    no_autofit(ptf)
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ptf.margin_left = Inches(0.08)
    sec = section_name or "SECTION"
    ptf.text = f"§ {str(section_num).zfill(2)} / {sec.upper()}"
    _style(ptf, BRAND if not dark else WHITE, MONO_LABEL_SIZE, bold=True, mono=True)

    title_x = CHROME_PAD + 1.7
    title_text = section_title or ""
    if title_text:
        ttb = slide.shapes.add_textbox(Inches(title_x), Inches(y - 0.02), Inches(6.5), Inches(0.35))
        ttf = ttb.text_frame
        no_autofit(ttf)
        ttf.text = title_text
        if section_cn:
            p = ttf.add_paragraph()
            p.text = f"· {section_cn}"
            for r in p.runs:
                r.font.name = FONT_CN
                r.font.size = Pt(11)
                r.font.color.rgb = INK_3 if not dark else BRAND_GLOW
        _style(ttf, INK_2 if not dark else WHITE, 12, bold=True, align=PP_ALIGN.LEFT)

    if rhs_meta:
        rtb = slide.shapes.add_textbox(Inches(9.2), Inches(y), Inches(3.5), Inches(0.3))
        rtf = rtb.text_frame
        no_autofit(rtf)
        rtf.text = rhs_meta
        rtf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        _style(rtf, INK_4 if not dark else BRAND_GLOW, MONO_LABEL_SIZE, mono=True, align=PP_ALIGN.RIGHT)

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(CHROME_PAD), Inches(y + 0.38),
        Inches(SLIDE_W - 2 * CHROME_PAD), Inches(0.01),
    )
    fill(rule, LINE if not dark else BRAND_2)
    return pill


def add_slide_footer(slide, left_text="© MetaStellar", sheet_num=1, total=20, dark=False):
    """Bottom chrome — brand line left, sheet id right."""
    y = SLIDE_H - 0.55
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(CHROME_PAD), Inches(y),
        Inches(SLIDE_W - 2 * CHROME_PAD), Inches(0.01),
    )
    fill(rule, LINE if not dark else BRAND_2)
    ltb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(y + 0.1), Inches(6.0), Inches(0.28))
    ltf = ltb.text_frame
    no_autofit(ltf)
    ltf.text = left_text
    _style(ltf, INK_4 if not dark else BRAND_GLOW, 9, mono=True)
    rtb = slide.shapes.add_textbox(Inches(10.5), Inches(y + 0.1), Inches(2.2), Inches(0.28))
    rtf = rtb.text_frame
    no_autofit(rtf)
    rtf.text = f"{str(sheet_num).zfill(2)} / {str(total).zfill(2)}"
    rtf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _style(rtf, INK_4 if not dark else BRAND_GLOW, 9, mono=True, align=PP_ALIGN.RIGHT)


def kpi_card(slide, x, y, w, h, label, value, unit=None, desc=None, badge=None, featured=False):
    """KPI tile — core data UI component (1 featured = BRAND_DEEP anchor)."""
    bg = BRAND_DEEP if featured else PAPER_2
    tc = WHITE if featured else INK
    sub_c = BRAND_GLOW if featured else INK_3
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(card, bg)
    card.line.color.rgb = BRAND_DEEP if featured else LINE
    card.line.width = Pt(1)

    ltb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.2), Inches(w - 0.44), Inches(0.3))
    ltf = ltb.text_frame
    no_autofit(ltf)
    ltf.text = label.upper()
    _style(ltf, sub_c, MONO_LABEL_SIZE, bold=True, mono=True)

    if badge:
        btw = min(1.1, w * 0.38)
        badge_sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + w - btw - 0.22), Inches(y + 0.18), Inches(btw), Inches(0.26),
        )
        fill(badge_sh, BRAND_3 if featured else BRAND_TINT)
        badge_sh.line.color.rgb = BRAND_2 if featured else BRAND_PALE
        badge_sh.line.width = Pt(1)
        btf = badge_sh.text_frame
        no_autofit(btf)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        btf.text = badge
        _style(btf, WHITE if featured else BRAND, 9, bold=True, mono=True, align=PP_ALIGN.CENTER)

    vtb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + h - 1.35), Inches(w - 0.44), Inches(0.95))
    vtf = vtb.text_frame
    no_autofit(vtf)
    vtf.text = str(value)
    _style(vtf, tc, KPI_VALUE_SIZE, bold=True, align=PP_ALIGN.LEFT)
    if unit:
        p = vtf.add_paragraph()
        p.text = unit
        for r in p.runs:
            r.font.name = FONT_MONO
            r.font.size = Pt(14)
            r.font.color.rgb = sub_c

    if desc:
        dtb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + h - 0.42), Inches(w - 0.44), Inches(0.32))
        dtf = dtb.text_frame
        no_autofit(dtf)
        dtf.text = desc
        _style(dtf, sub_c, 11, bold=False)
    return card


def new_deck():
    """16:9 presentation 16:9 widescreen geometry."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def title_only_slide(prs):
    """Reproduce the 'Title Only' layout: blank slide; place shapes yourself."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _place_logo(slide, logo_path, left, top, height=0.48):
    if logo_path:
        slide.shapes.add_picture(logo_path, Inches(left), Inches(top), height=Inches(height))


def presentation_title_slide(prs, title, presenter=None, logo_path=None):
    """Main title — BRAND_DEEP field, MetaStellar mark, white headline."""
    slide = title_only_slide(prs)
    _slide_bg(slide, BRAND_DEEP)
    add_wordmark(slide, CHROME_PAD, 0.35, dark=True, compact=True)
    tb = slide.shapes.add_textbox(Inches(1.1), Inches(2.55), Inches(9.5), Inches(1.15))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    _style(tf, WHITE, 36, bold=True, align=PP_ALIGN.LEFT)
    if presenter:
        ptb = slide.shapes.add_textbox(Inches(1.1), Inches(3.85), Inches(8.0), Inches(0.55))
        ptf = ptb.text_frame
        no_autofit(ptf)
        ptf.text = presenter
        _style(ptf, BRAND_GLOW, 18, bold=False, align=PP_ALIGN.LEFT)
    _place_logo(slide, logo_path, 11.35, 0.38, height=0.5)
    add_slide_footer(slide, sheet_num=1, total=20, dark=True)
    return slide


def agenda_slide(prs, topics, title="Agenda", logo_path=None, columns=1, sheet_num=6, total=20):
    """Agenda — left hero + right numbered rows with § pills."""
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_wordmark(slide, CHROME_PAD, 0.35, compact=True)
    tb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(1.2), Inches(4.2), Inches(1.0))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    _style(tf, INK, 28, bold=True, align=PP_ALIGN.LEFT)
    n = len(topics)
    row_y0 = 1.15
    row_gap = 0.72
    rx = 5.2
    for i, topic in enumerate(topics):
        y = row_y0 + i * row_gap
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(rx), Inches(y), Inches(0.55), Inches(0.28),
        )
        fill(pill, BRAND_TINT)
        pill.line.color.rgb = BRAND_PALE
        pill.line.width = Pt(1)
        ptf = pill.text_frame
        no_autofit(ptf)
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ptf.text = f"§ {str(i + 1).zfill(2)}"
        _style(ptf, BRAND, 9, bold=True, mono=True, align=PP_ALIGN.CENTER)
        label = slide.shapes.add_textbox(Inches(rx + 0.68), Inches(y - 0.02), Inches(6.8), Inches(0.38))
        ltf = label.text_frame
        no_autofit(ltf)
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ltf.text = topic
        _style(ltf, INK_2, 14, bold=False, align=PP_ALIGN.LEFT)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.42)
    add_slide_footer(slide, sheet_num=sheet_num, total=total)
    return slide


def section_title_slide(prs, title, image_path=None, logo_path=None, section_num=7):
    """Section divider — BRAND_DEEP + § number + kicker."""
    slide = title_only_slide(prs)
    _slide_bg(slide, BRAND_DEEP)
    add_wordmark(slide, CHROME_PAD, 0.35, dark=True, compact=True)
    num_tb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(2.0), Inches(5.5), Inches(1.8))
    ntf = num_tb.text_frame
    no_autofit(ntf)
    ntf.text = str(section_num).zfill(2)
    _style(ntf, BRAND_GLOW, 72, bold=True, align=PP_ALIGN.LEFT)
    tb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(3.65), Inches(10.5), Inches(1.0))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    _style(tf, WHITE, 32, bold=True, align=PP_ALIGN.LEFT)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(7.5), Inches(1.65),
                                 width=Inches(5.0), height=Inches(4.35))
    add_slide_footer(slide, sheet_num=section_num, total=20, dark=True)
    _place_logo(slide, logo_path, 0.45, 6.85, height=0.38)
    return slide


def _place_image_or_placeholder(slide, image_path, x, y, w, h, label="Image"):
    """Use a real provided image when available; otherwise leave a gray placeholder."""
    if image_path:
        slide.shapes.add_picture(image_path, Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
        return None
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    fill(ph, GRAY_2)
    ph.line.color.rgb = GRAY_3
    ph.line.width = Pt(1)
    tf = ph.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = label
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    _style(tf, BLUE_800, 11, bold=False, align=PP_ALIGN.CENTER)
    return ph


def left_image_icon_rows_slide(
    prs,
    title,
    rows,
    punchline=None,
    img_path=None,
    logo_path=None,
    img_placeholder="Image",
):
    """Left hero image + accent icon tiles + gray statement rows (see layout-library #14)."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    img_w = 4.15
    _place_image_or_placeholder(slide, img_path, 0, 0, img_w, SLIDE_H, label=img_placeholder)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.4)
    rx = img_w + 0.35
    rw = SLIDE_W - rx - 0.35
    tb = slide.shapes.add_textbox(Inches(rx), Inches(0.32), Inches(rw), Inches(0.7))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    _style(tf, BLUE_800, 26, bold=True, align=PP_ALIGN.LEFT)
    row_h = 0.52
    gap = 0.14
    y0 = 1.25
    tile = 0.48
    for i, row in enumerate(rows):
        y = y0 + i * (row_h + gap)
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(y),
                                    Inches(tile), Inches(row_h))
        fill(sq, ACCENT)
        icon = row.get("icon_path")
        if icon:
            slide.shapes.add_picture(icon, Inches(rx + 0.06), Inches(y + 0.06),
                                     width=Inches(tile - 0.12), height=Inches(row_h - 0.12))
        bar_x = rx + tile + 0.12
        bar_w = rw - tile - 0.12
        box(slide, bar_x, y, bar_w, row_h, row.get("text", ""), GRAY_1,
            BLUE_800, size=13, bold=False, align=PP_ALIGN.LEFT)
    if punchline:
        ptb = slide.shapes.add_textbox(Inches(rx), Inches(5.85), Inches(rw), Inches(1.2))
        ptf = ptb.text_frame
        no_autofit(ptf)
        ptf.word_wrap = True
        ptf.text = punchline
        _style(ptf, BLUE_800, 22, bold=True, align=PP_ALIGN.LEFT)
    return slide


def left_image_tiered_list_slide(
    prs,
    title,
    message,
    categories,
    img_path=None,
    logo_path=None,
    img_placeholder="Image",
):
    """Left hero + overlapping navy message bar + accent/gray category rows (layout-library #16)."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    img_w = 5.05
    img_top = 1.35
    img_h = 5.95
    _place_image_or_placeholder(slide, img_path, 0, img_top, img_w, img_h, label=img_placeholder)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    rx = img_w + 0.25
    rw = SLIDE_W - rx - 0.35
    tb = slide.shapes.add_textbox(Inches(rx), Inches(0.28), Inches(rw), Inches(0.65))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    _style(tf, BLUE_800, 24, bold=True, align=PP_ALIGN.LEFT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(0.88),
                                 Inches(10.2), Inches(0.78))
    fill(bar, BLUE_800)
    btf = bar.text_frame
    no_autofit(btf)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.margin_left = Inches(0.25)
    btf.margin_right = Inches(0.25)
    btf.text = message
    _style(btf, WHITE, 16, bold=True, align=PP_ALIGN.LEFT)
    y = 1.45
    head_w = 1.05
    gap = 0.12
    for cat in categories:
        body_h = max(0.95, 0.28 + 0.22 * len(cat.get("bullets", [])))
        box(slide, rx, y, head_w, body_h, cat.get("title", ""), ACCENT,
            BLUE_800, size=13, bold=True, align=PP_ALIGN.LEFT)
        bullet_box(slide, rx + head_w + gap, y, rw - head_w - gap, body_h,
                   cat.get("bullets", []), shade=GRAY_1, size=12)
        y += body_h + gap
    return slide


def migration_timeline_circles_slide(
    prs,
    title,
    subtitle=None,
    steps=None,
    logo_path=None,
):
    """Wave timeline with alternating circles + accent callout markers (layout-library #15)."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    steps = steps or []
    tb = slide.shapes.add_textbox(Inches(7.8), Inches(0.28), Inches(5.2), Inches(0.95))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _style(tf, BLUE_800, 26, bold=True, align=PP_ALIGN.RIGHT)
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.alignment = PP_ALIGN.RIGHT
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(18)
            r.font.color.rgb = BLUE_800
    n = len(steps)
    if not n:
        return slide
    d = 1.75
    x0 = 0.75
    span = SLIDE_W - 1.5
    gap = (span - d * n) / (n - 1) if n > 1 else 0
    cy = 3.55
    arrow(slide, x0, cy, x0 + span, cy, color=BLUE_800, width=1.2, dashed=True)
    for i, step in enumerate(steps):
        x = x0 + i * (d + gap)
        col = BLUE_800 if i % 2 == 0 else RAMP[2]
        tc = text_on(col)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(cy - d / 2),
                                      Inches(d), Inches(d))
        fill(circ, col)
        ctf = circ.text_frame
        no_autofit(ctf)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.word_wrap = True
        ctf.text = step.get("label", "")
        for p in ctf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        _style(ctf, tc, 11, bold=True, align=PP_ALIGN.CENTER)
        callout = step.get("callout")
        if callout:
            above = step.get("callout_pos", "above") != "below"
            my = cy - d / 2 - 0.55 if above else cy + d / 2 + 0.18
            dot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + d / 2 - 0.03),
                                         Inches(my - 0.18 if above else cy + d / 2),
                                         Inches(0.06), Inches(0.28 if above else 0.22))
            fill(dot, ACCENT)
            label(slide, x + d / 2 - 1.35, my - 0.15 if above else my,
                  2.7, callout, size=10, bold=False, align=PP_ALIGN.CENTER)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def no_autofit(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def corner_accent(slide):
    """Top-right brand tile (replaces indigo accent block)."""
    a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.35), Inches(0.0),
                               Inches(0.78), Inches(0.34))
    fill(a, BRAND)
    return a


def _style_runs_legacy(tf, color, size_pt, bold=False, font=FONT, align=None):
    """Alias kept for internal helpers that predate mono support."""
    _style(tf, color, size_pt, bold=bold, font=font, align=align)


def add_title(slide, text, subtitle=None, align=PP_ALIGN.RIGHT, size=27):
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(0.28), Inches(12.1), Inches(0.95))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = text
    tf.paragraphs[0].alignment = align
    _style(tf, BLUE_800, size, bold=True)
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.alignment = align
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(13)
            r.font.color.rgb = ACCENT_LIGHT
    return tb


def add_message_bar(slide, text, fill_color=ACCENT_LIGHT):
    """The key message bar. 20pt Segoe UI; use theme palette fills only."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.43),
                                 Inches(9.8), Inches(0.84))
    fill(bar, fill_color)
    tf = bar.text_frame
    no_autofit(tf)
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = text
    txt = text_on(fill_color)
    _style(tf, txt, MSG_BAR_SIZE, bold=True)
    return bar


def box(slide, x, y, w, h, text, fill_color, text_color=None, size=BODY_SIZE, bold=True,
        align=PP_ALIGN.CENTER, sub=None, sub_size=12):
    """A rectangle with embedded text. text_color defaults to text_on(fill)."""
    if text_color is None:
        text_color = text_on(fill_color)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    fill(sh, fill_color)
    tf = sh.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.text = text
    p0 = tf.paragraphs[0]
    p0.alignment = align
    for r in p0.runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = text_color
    if sub:
        p = tf.add_paragraph()
        p.text = sub
        p.alignment = align
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(sub_size)
            r.font.color.rgb = text_color
    return sh


def bullet_box(slide, x, y, w, h, bullets, shade=GRAY_1, size=BODY_SIZE, heading=None, text_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    fill(sh, shade)
    tc = text_color or BLUE_800
    tf = sh.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_right = Inches(0.14)
    first = True
    if heading:
        tf.text = heading
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        for r in tf.paragraphs[0].runs:
            r.font.name = FONT
            r.font.size = Pt(size + 2)
            r.font.bold = True
            r.font.color.rgb = BLUE_800
        first = False
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "\u2022 " + b
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = tc
    return sh


def label(slide, x, y, w, text, color=BLUE_800, size=11, bold=True,
          align=PP_ALIGN.LEFT, mono=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    tf = tb.text_frame
    no_autofit(tf)
    tf.margin_left = Inches(0.02)
    tf.margin_top = Inches(0)
    tf.text = text
    tf.paragraphs[0].alignment = align
    for r in tf.paragraphs[0].runs:
        r.font.name = FONT_MONO if mono else FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def arrow(slide, x1, y1, x2, y2, color=BLUE_800, width=2.25, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),
                             {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if dashed:
        ln.insert(0, ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return c


def dashed_container(slide, x, y, w, h, color=SLATE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = Pt(1.5)
    ln = sh.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    sh.shadow.inherit = False
    return sh


def add_activation_timeline(
    slide,
    checkpoints,
    top_items=None,
    bottom_items=None,
    x0=0.55,
    x1=12.8,
    y_line=3.95,
):
    """Draw a 2-lane activation timeline similar to GTM launch plans.

    checkpoints: [{"x": 2.4, "label": "Mar 5"}, ...]
    top_items / bottom_items:
      [{"x": 2.6, "title": "Milestone", "note": "optional", "w": 2.0, "color": ACCENT}, ...]
    """
    top_items = top_items or []
    bottom_items = bottom_items or []

    # Main axis
    arrow(slide, x0, y_line, x1, y_line, color=BLUE_800, width=2.2)

    # Lane labels
    box(slide, 0.0, y_line - 1.55, 0.42, 1.42, "CONTENT", BLUE_800, WHITE, size=10)
    box(slide, 0.0, y_line + 0.15, 0.42, 1.42, "PROMO", BLUE_800, WHITE, size=10)

    # Checkpoints on axis
    for cp in checkpoints:
        x = cp["x"]
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.08), Inches(y_line - 0.08),
                                     Inches(0.16), Inches(0.16))
        fill(dot, WHITE)
        dot.line.color.rgb = BLUE_800
        dot.line.width = Pt(1.5)
        if cp.get("label"):
            label(slide, x - 0.35, y_line - 0.42, 0.8, cp["label"], color=BLUE_800,
                  size=10, bold=False, align=PP_ALIGN.CENTER)

    # Activity cards above axis
    for i, item in enumerate(top_items):
        w = item.get("w", 2.0)
        h = 0.5 if not item.get("note") else 0.82
        y = y_line - 1.2 - (0.62 * (i % 2))
        fc = item.get("color", ACCENT)
        tc = text_on(fc)
        box(slide, item["x"] - w / 2, y, w, h, item["title"], fc, tc, size=11,
            align=PP_ALIGN.LEFT, sub=item.get("note"), sub_size=10)
        arrow(slide, item["x"], y + h, item["x"], y_line - 0.12, color=GRAY_3, width=1.2)

    # Activity cards below axis
    for i, item in enumerate(bottom_items):
        w = item.get("w", 2.0)
        h = 0.5 if not item.get("note") else 0.82
        y = y_line + 0.26 + (0.62 * (i % 2))
        fc = item.get("color", GRAY_1)
        tc = text_on(fc)
        box(slide, item["x"] - w / 2, y, w, h, item["title"], fc, tc, size=11,
            align=PP_ALIGN.LEFT, sub=item.get("note"), sub_size=10)
        arrow(slide, item["x"], y, item["x"], y_line + 0.12, color=GRAY_3, width=1.2)


def timeline_template_slide(prs, title, message, checkpoints, top_items=None, bottom_items=None):
    """One-call template: title + message bar + 2-lane timeline."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    add_title(slide, title, align=PP_ALIGN.RIGHT, size=TITLE_SIZE + 2)
    add_message_bar(slide, message, fill_color=ACCENT_LIGHT)
    add_activation_timeline(slide, checkpoints, top_items=top_items, bottom_items=bottom_items)
    return slide


def add_cards_row(slide, cards, top=2.5, bottom=6.95, gap=0.4, left_margin=0.45,
                  header="accent", img_ratio=2.35):
    """cards = [{"title", "bullets":[...], "img": path|None}].
    header = 'accent' (indigo accent bar) or 'ramp' (graded blue ramp).
    Lays out N equal cards: header bar + optional image banner + gray bullet box.
    """
    from PIL import Image
    n = len(cards)
    total_w = SLIDE_W - 2 * left_margin
    w = (total_w - gap * (n - 1)) / n
    head_h = 0.5
    img_top = top + head_h + 0.08
    img_h = w / img_ratio
    box_top = img_top + img_h + 0.08
    for i, c in enumerate(cards):
        x = left_margin + i * (w + gap)
        if header == "accent":
            hf, ht = ACCENT, BLUE_800
        else:
            hf, ht = RAMP[min(i, 3)], text_on(RAMP[min(i, 3)])
        box(slide, x, top, w, head_h, c["title"], hf, ht, size=BODY_SIZE,
            align=PP_ALIGN.LEFT)
        bt = box_top
        img = c.get("img")
        if img:
            # crop to a clean banner, then place with locked aspect ratio
            im = Image.open(img)
            iw, ih = im.size
            crop_h = min(ih, int(iw / img_ratio))
            cropped = img + ".banner.png"
            im.crop((0, 0, iw, crop_h)).save(cropped)
            slide.shapes.add_picture(cropped, Inches(x), Inches(img_top),
                                     width=Inches(w), height=Inches(img_h))
        else:
            bt = top + head_h
        bullet_box(slide, x, bt, w, top + (bottom - top) - bt,
                   c.get("bullets", []))


def footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.32))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = text
    for r in tf.paragraphs[0].runs:
        r.font.name = FONT
        r.font.size = Pt(9)
        r.font.color.rgb = BLUE_800
    return tb


# ---- Marketing agency / client pitch layouts (premium) ----

def agency_cover_slide(prs, headline=None, client=None, subtitle=None, date=None, logo_path=None,
                       metrics=None, sheet_num=1, total=DEFAULT_TOTAL):
    """Opener — left hero title + right BRAND_DEEP metric panel (SEO/GEO defaults)."""
    headline = headline or SEO_GEO_TAGLINE
    client = client if client is not None else ""
    subtitle = subtitle or (
        "能见度的新科学 —— 让品牌在人类和 AI 的每一次搜索中都能被找到、被引用、被信任。"
    )
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_wordmark(slide, CHROME_PAD, 0.35, compact=True)
    tb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(1.85), Inches(7.2), Inches(2.2))
    tf = tb.text_frame
    no_autofit(tf)
    tf.word_wrap = True
    tf.text = headline
    _style(tf, INK, 40, bold=True, align=PP_ALIGN.LEFT)
    if client:
        ctb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(4.15), Inches(7.0), Inches(0.55))
        ctf = ctb.text_frame
        no_autofit(ctf)
        ctf.text = client
        _style(ctf, BRAND, 20, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        y_sub = 4.75 if client else 4.15
        stb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(y_sub), Inches(7.0), Inches(0.9))
        stf = stb.text_frame
        no_autofit(stf)
        stf.word_wrap = True
        stf.text = subtitle
        _style(stf, INK_3, 13, bold=False, align=PP_ALIGN.LEFT)
    panel_x = 8.35
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(panel_x), Inches(0), Inches(SLIDE_W - panel_x), Inches(SLIDE_H),
    )
    fill(panel, BRAND_DEEP)
    metrics = metrics or DEFAULT_SEO_METRICS
    my = 0.85
    mh = 1.35
    for m in metrics[:4]:
        kpi_card(slide, panel_x + 0.35, my, SLIDE_W - panel_x - 0.7, mh,
                 m.get("label", ""), m.get("value", ""), unit=m.get("unit"),
                 desc=m.get("desc"), badge=m.get("badge"), featured=False)
        my += mh + 0.22
    if date:
        add_slide_footer(slide, left_text=date, sheet_num=sheet_num, total=total)
    else:
        add_slide_footer(slide, sheet_num=sheet_num, total=total)
    _place_logo(slide, logo_path, panel_x + 0.35, 0.35, height=0.42)
    return slide


def agency_section_slide(prs, section_num, title, subtitle=None, logo_path=None, sheet_num=7, total=20):
    """Section divider — BRAND_DEEP, oversized § index, kicker pill."""
    slide = title_only_slide(prs)
    _slide_bg(slide, BRAND_DEEP)
    add_wordmark(slide, CHROME_PAD, 0.35, dark=True, compact=True)
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(CHROME_PAD), Inches(1.35), Inches(1.8), Inches(0.3),
    )
    fill(pill, BRAND)
    ptf = pill.text_frame
    no_autofit(ptf)
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ptf.margin_left = Inches(0.1)
    ptf.text = f"§ {str(section_num).zfill(2)}"
    _style(ptf, WHITE, MONO_LABEL_SIZE, bold=True, mono=True)
    num_tb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(2.0), Inches(5.0), Inches(1.6))
    ntf = num_tb.text_frame
    no_autofit(ntf)
    ntf.text = str(section_num).zfill(2)
    _style(ntf, BRAND_GLOW, 88, bold=True, align=PP_ALIGN.LEFT)
    tb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(3.55), Inches(11.0), Inches(1.0))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    _style(tf, WHITE, 34, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        stb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(4.55), Inches(11.0), Inches(0.55))
        stf = stb.text_frame
        no_autofit(stf)
        stf.text = subtitle
        _style(stf, BRAND_GLOW, 16, bold=False, align=PP_ALIGN.LEFT)
    add_slide_footer(slide, sheet_num=sheet_num, total=total, dark=True)
    _place_logo(slide, logo_path, 11.35, 6.55, height=0.38)
    return slide


def big_idea_slide(prs, statement=None, caption=None, logo_path=None, section_num=8, sheet_num=8, total=DEFAULT_TOTAL):
    """Statement / manifesto — single quote + § chrome."""
    statement = statement or SEO_GEO_MANIFESTO
    caption = caption or "如果无法被找到，就等于不存在 —— 能见度，就是新的存在权。"
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_slide_header(slide, section_num, "INSIGHT", section_title="Statement")
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.3), Inches(2.4))
    tf = tb.text_frame
    no_autofit(tf)
    tf.word_wrap = True
    tf.text = statement
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _style(tf, INK, 36, bold=True, align=PP_ALIGN.CENTER)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.4), Inches(4.75),
                                  Inches(2.5), Inches(0.04))
    fill(rule, BRAND)
    if caption:
        ctb = slide.shapes.add_textbox(Inches(1.5), Inches(5.05), Inches(10.3), Inches(0.55))
        ctf = ctb.text_frame
        no_autofit(ctf)
        ctf.text = caption
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _style(ctf, INK_3, 15, bold=False, align=PP_ALIGN.CENTER)
    add_slide_footer(slide, sheet_num=sheet_num, total=total)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def metrics_3up_slide(prs, title, metrics, message=None, logo_path=None,
                      section_num=9, sheet_num=9, total=20):
    """Data / KPIs — up to 4 KPI cards (last featured on BRAND_DEEP).

    metrics: [{"value": "38%", "label": "Reach uplift", "note": "vs baseline", "badge": "↑ 18%"}, ...]
    """
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_slide_header(slide, section_num, "DATA", section_title=title, rhs_meta="KPIs")
    if message:
        mbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(CHROME_PAD), Inches(1.15),
                                      Inches(10.5), Inches(0.55))
        fill(mbar, BRAND_TINT)
        mbar.line.color.rgb = BRAND_PALE
        mbar.line.width = Pt(1)
        mtf = mbar.text_frame
        no_autofit(mtf)
        mtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        mtf.margin_left = Inches(0.2)
        mtf.text = message
        _style(mtf, INK_2, 13, bold=False)
    n = min(len(metrics), 4)
    if not n:
        add_slide_footer(slide, sheet_num=sheet_num, total=total)
        return slide
    gap = 0.35
    left = CHROME_PAD
    total_w = SLIDE_W - 2 * CHROME_PAD
    w = (total_w - gap * (n - 1)) / n
    top = 2.05 if message else 1.55
    h = 4.2
    for i, m in enumerate(metrics[:4]):
        x = left + i * (w + gap)
        featured = i == n - 1 and n >= 3
        kpi_card(slide, x, top, w, h,
                 m.get("label", "METRIC"),
                 m.get("value", ""),
                 unit=m.get("unit"),
                 desc=m.get("note"),
                 badge=m.get("badge"),
                 featured=featured)
    add_slide_footer(slide, sheet_num=sheet_num, total=total)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def challenge_solution_slide(prs, title, challenge_bullets, solution_bullets, logo_path=None):
    """Split challenge vs. agency approach — proposal staple."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    add_title(slide, title, align=PP_ALIGN.LEFT, size=26)
    col_w = 5.85
    gap = 0.55
    x0 = 0.55
    x1 = x0 + col_w + gap
    top = 1.35
    head_h = 0.52
    body_h = 5.15
    box(slide, x0, top, col_w, head_h, "Challenge", GRAY_800, WHITE, size=14, bold=True)
    bullet_box(slide, x0, top + head_h + 0.08, col_w, body_h, challenge_bullets[:5],
               shade=GRAY_100, size=13)
    box(slide, x1, top, col_w, head_h, "Our approach", BLUE_600, WHITE, size=14, bold=True)
    bullet_box(slide, x1, top + head_h + 0.08, col_w, body_h, solution_bullets[:5],
               shade=BLUE_50, size=13)
    _place_logo(slide, logo_path, 11.35, 6.55, height=0.38)
    return slide


def pillar_strategy_slide(prs, title, pillars, message=None, logo_path=None):
    """3–4 strategic pillars — agency strategy / recommendation slide.

    pillars: [{"title": "...", "bullets": ["...", ...]}, ...]
    """
    slide = title_only_slide(prs)
    corner_accent(slide)
    add_title(slide, title, align=PP_ALIGN.LEFT, size=26)
    if message:
        add_message_bar(slide, message, fill_color=ACCENT_LIGHT)
    cards = [{"title": p.get("title", ""), "bullets": p.get("bullets", [])} for p in pillars[:4]]
    add_cards_row(slide, cards, top=2.2 if message else 1.55, bottom=6.85,
                  header="accent", img_ratio=3.5)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def campaign_timeline_slide(prs, title, message, checkpoints, top_items=None, bottom_items=None):
    """Agency campaign / rollout timeline — wraps timeline_template_slide."""
    return timeline_template_slide(prs, title, message, checkpoints,
                                   top_items=top_items, bottom_items=bottom_items)


CLOSING_FOOTER = "Thank you"


def contact_slide(prs, headline=None, subtitle=None,
                  contacts=None, sheet_num=24, total=DEFAULT_TOTAL, logo_path=None):
    """Contact close — BRAND_DEEP + CTA + glass-style contact rows."""
    headline = headline or SEO_GEO_CONTACT
    slide = title_only_slide(prs)
    _slide_bg(slide, BRAND_DEEP)
    add_wordmark(slide, CHROME_PAD, 0.35, dark=True, compact=True)
    tb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(2.2), Inches(6.5), Inches(1.8))
    tf = tb.text_frame
    no_autofit(tf)
    tf.word_wrap = True
    tf.text = headline
    _style(tf, WHITE, 36, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        stb = slide.shapes.add_textbox(Inches(CHROME_PAD), Inches(4.05), Inches(6.0), Inches(0.55))
        stf = stb.text_frame
        no_autofit(stf)
        stf.text = subtitle
        _style(stf, BRAND_GLOW, 15, bold=False, align=PP_ALIGN.LEFT)
    cta = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(CHROME_PAD), Inches(4.85), Inches(2.4), Inches(0.48),
    )
    fill(cta, BRAND_3)
    ctf = cta.text_frame
    no_autofit(ctf)
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.text = "Get in touch →"
    _style(ctf, WHITE, 14, bold=True, align=PP_ALIGN.CENTER)
    contacts = contacts or [
        ("Studio", "Shanghai · Remote"),
        ("Email", "hello@metastellar.com"),
        ("Phone", "+86 · · ·"),
        ("Follow", "@metastellar"),
    ]
    rx = 7.8
    ry = 1.65
    for label, value in contacts[:4]:
        row = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(rx), Inches(ry), Inches(4.8), Inches(0.72),
        )
        row.fill.solid()
        row.fill.fore_color.rgb = BRAND
        row.fill.transparency = 0.92
        row.line.color.rgb = BRAND_GLOW
        row.line.width = Pt(1)
        ltb = slide.shapes.add_textbox(Inches(rx + 0.2), Inches(ry + 0.12), Inches(4.4), Inches(0.5))
        ltf = ltb.text_frame
        no_autofit(ltf)
        ltf.text = f"{label.upper()}  ·  {value}"
        _style(ltf, WHITE, 11, mono=True)
        ry += 0.88
    add_slide_footer(slide, sheet_num=sheet_num, total=total, dark=True)
    _place_logo(slide, logo_path, rx, 0.35, height=0.42)
    return slide


def closing_slide(prs, logo_path=None, tagline=None):
    """Thank-you slide — SEO/GEO contact close."""
    return contact_slide(
        prs, headline=tagline or SEO_GEO_CONTACT, subtitle="Q&A",
        sheet_num=DEFAULT_TOTAL, total=DEFAULT_TOTAL, logo_path=logo_path,
    )


# ---- MetaStellar SEO/GEO specialty slides (24-page deck) ----

def _engine_card(slide, x, y, w, h, name, owner, stat, desc, tag=None, featured=False):
    bg = BRAND_DEEP if featured else PAPER_2
    tc = WHITE if featured else INK
    sub = BRAND_GLOW if featured else INK_3
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(card, bg)
    card.line.color.rgb = BRAND_DEEP if featured else LINE
    card.line.width = Pt(1)
    if tag:
        label(slide, x + 0.15, y + 0.15, w - 0.3, tag.upper(), color=sub, size=8, bold=True, align=PP_ALIGN.LEFT)
    label(slide, x + 0.15, y + 0.42, w - 0.3, name, color=tc, size=14, bold=True)
    label(slide, x + 0.15, y + 0.72, w - 0.3, owner, color=sub, size=10, bold=False)
    vtb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.05), Inches(w - 0.3), Inches(0.7))
    vtf = vtb.text_frame
    no_autofit(vtf)
    vtf.text = stat
    _style(vtf, tc, 28, bold=True)
    if desc:
        label(slide, x + 0.15, y + h - 0.55, w - 0.3, desc, color=sub, size=9, bold=False)


def ai_search_landscape_slide(prs, title="Four engines, one new frontier.",
                              engines=None, logo_path=None, sheet_num=10, total=DEFAULT_TOTAL):
    """Slide 10 — AI Search Landscape: 4 engine comparison cards."""
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_slide_header(slide, 10, "LANDSCAPE", section_title=title, rhs_meta="AI Search")
    engines = engines or [
        {"name": "ChatGPT", "owner": "OpenAI", "stat": "3.7B", "desc": "queries / week", "tag": "featured", "featured": True},
        {"name": "Perplexity", "owner": "Perplexity AI", "stat": "100M+", "desc": "monthly users"},
        {"name": "Google AIO", "owner": "Google", "stat": "↓34.5%", "desc": "organic CTR impact", "tag": "SERP"},
        {"name": "Gemini", "owner": "Google / MS", "stat": "58%", "desc": "zero-click searches", "tag": "Copilot"},
    ]
    gap, left, top, h = 0.28, CHROME_PAD, 1.55, 4.85
    n = min(len(engines), 4)
    w = (SLIDE_W - 2 * CHROME_PAD - gap * (n - 1)) / n
    for i, e in enumerate(engines[:4]):
        _engine_card(slide, left + i * (w + gap), top, w, h,
                     e["name"], e.get("owner", ""), e.get("stat", ""),
                     e.get("desc", ""), tag=e.get("tag"), featured=e.get("featured", False))
    add_slide_footer(slide, sheet_num=sheet_num, total=total)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def geo_explained_slide(prs, title="When AI answers, who gets quoted?",
                        seo_steps=None, geo_steps=None, shift_note=None,
                        logo_path=None, sheet_num=11, total=DEFAULT_TOTAL):
    """Slide 11 — GEO Explained: Traditional SEO vs Generative Search dual flow."""
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_slide_header(slide, 11, "GEO", section_title=title)
    seo_steps = seo_steps or ["Crawl & index", "Rank keywords", "Earn backlinks", "Win SERP clicks"]
    geo_steps = geo_steps or ["Train on sources", "Answer synthesis", "Cite brands", "Win AI citations"]
    col_w, gap, top, body_h = 5.85, 0.45, 1.55, 3.85
    x0, x1 = CHROME_PAD, CHROME_PAD + col_w + gap
    box(slide, x0, top, col_w, 0.45, "Traditional SEO", PAPER_3, INK_2, size=12, bold=True)
    bullet_box(slide, x0, top + 0.52, col_w, body_h, seo_steps, shade=PAPER_2, size=13)
    box(slide, x1, top, col_w, 0.45, "Generative Search (GEO)", BRAND_DEEP, WHITE, size=12, bold=True)
    bullet_box(slide, x1, top + 0.52, col_w, body_h, geo_steps, shade=BRAND_TINT, size=13, text_color=BRAND)
    note = shift_note or "The Shift: 过去比拼排名，现在比拼被引用。"
    nbar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(CHROME_PAD), Inches(5.65),
                                  Inches(SLIDE_W - 2 * CHROME_PAD), Inches(0.55))
    fill(nbar, BRAND_TINT)
    nbar.line.color.rgb = BRAND_PALE
    nbar.line.width = Pt(1)
    ntf = nbar.text_frame
    no_autofit(ntf)
    ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ntf.margin_left = Inches(0.2)
    ntf.text = note
    _style(ntf, BRAND, 12, bold=True)
    add_slide_footer(slide, sheet_num=sheet_num, total=total)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def seo_vs_geo_slide(prs, title="Same goal. Different rulebook.",
                     rows=None, logo_path=None, sheet_num=12, total=DEFAULT_TOTAL):
    """Slide 12 — SEO vs GEO comparison table (6 dimensions)."""
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_slide_header(slide, 12, "COMPARE", section_title=title)
    rows = rows or [
        ("优化对象", "Keywords & pages", "Answers & citations"),
        ("核心指标", "Rank · CTR · Traffic", "Citation rate · Share of answer"),
        ("内容形态", "Landing pages · blogs", "Structured · quotable blocks"),
        ("技术信号", "Schema · Core Web Vitals", "Entity graph · source authority"),
        ("优化周期", "3–6 months", "4–8 weeks to first citations"),
        ("转化路径", "Click → site → convert", "Cite → trust → discover"),
    ]
    top, row_h, col_w = 1.55, 0.62, 4.2
    x_dim, x_seo, x_geo = CHROME_PAD, CHROME_PAD + 2.5, CHROME_PAD + 2.5 + col_w + 0.25
    box(slide, x_dim, top, 2.3, 0.4, "维度", PAPER_3, INK_2, size=11, bold=True)
    box(slide, x_seo, top, col_w, 0.4, "SEO", PAPER_3, INK_2, size=11, bold=True)
    box(slide, x_geo, top, col_w, 0.4, "GEO", BRAND, WHITE, size=11, bold=True)
    for i, (dim, seo, geo) in enumerate(rows[:6]):
        y = top + 0.48 + i * (row_h + 0.06)
        box(slide, x_dim, y, 2.3, row_h, dim, PAPER_2, INK_2, size=10, bold=True, align=PP_ALIGN.LEFT)
        box(slide, x_seo, y, col_w, row_h, seo, PAPER_2, INK_3, size=10, bold=False, align=PP_ALIGN.LEFT)
        box(slide, x_geo, y, col_w, row_h, geo, BRAND_TINT, BRAND, size=10, bold=False, align=PP_ALIGN.LEFT)
    add_slide_footer(slide, sheet_num=sheet_num, total=total)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def methodology_slide(prs, title="Five phases, one visibility loop.",
                      phases=None, logo_path=None, sheet_num=13, total=DEFAULT_TOTAL):
    """Slide 13 — Visibility Loop: Audit → Strategy → Build → Amplify → Measure."""
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_slide_header(slide, 13, "METHOD", section_title=title)
    phases = phases or [
        {"name": "Audit", "desc": "Technical + AI citation baseline", "items": ["SEO audit", "GEO snapshot", "Competitive map"]},
        {"name": "Strategy", "desc": "Keyword + answer architecture", "items": ["Topic clusters", "Citation targets", "Content briefs"]},
        {"name": "Build", "desc": "Content + structured data", "items": ["On-page SEO", "GEO content", "Schema deploy"]},
        {"name": "Amplify", "desc": "Campaign + distribution", "items": ["Paid media", "PR + links", "Social proof"]},
        {"name": "Measure", "desc": "Rank + citation dashboard", "items": ["Weekly KPIs", "AI monitoring", "Iteration"]},
    ]
    n = min(len(phases), 5)
    gap, left, top, w, h = 0.22, CHROME_PAD, 1.65, (SLIDE_W - 2 * CHROME_PAD - 0.22 * 4) / 5, 4.6
    track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top - 0.18),
                                   Inches(SLIDE_W - 2 * CHROME_PAD), Inches(0.04))
    fill(track, BRAND_PALE)
    for i, ph in enumerate(phases[:5]):
        x = left + i * (w + gap)
        featured = i == 2
        bg = BRAND_DEEP if featured else PAPER_2
        tc = WHITE if featured else INK
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(h))
        fill(card, bg)
        card.line.color.rgb = BRAND_DEEP if featured else LINE
        card.line.width = Pt(1)
        label(slide, x + 0.12, top + 0.15, w - 0.24, f"0{i + 1}", color=BRAND_GLOW if featured else BRAND, size=9, mono=True)
        label(slide, x + 0.12, top + 0.38, w - 0.24, ph.get("name", ""), color=tc, size=13, bold=True)
        label(slide, x + 0.12, top + 0.72, w - 0.24, ph.get("desc", ""), color=BRAND_GLOW if featured else INK_3, size=9, bold=False)
        items = ph.get("items", [])
        bullet_box(slide, x + 0.08, top + 1.15, w - 0.16, h - 1.3, items[:3],
                   shade=bg, size=9, text_color=WHITE if featured else INK_3)
    add_slide_footer(slide, sheet_num=sheet_num, total=total)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def services_4up_slide(prs, title="SEO · GEO · Campaign · Strategy",
                       services=None, logo_path=None, sheet_num=14, total=DEFAULT_TOTAL):
    """Slide 14 — 4 core services (GEO featured)."""
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_slide_header(slide, 14, "SERVICES", section_title=title)
    services = services or [
        {"num": "01", "title": "SEO", "cn": "搜索引擎优化",
         "desc": "Traditional search visibility",
         "subs": ["Technical Audit", "Keyword Strategy", "Backlink Building", "SERP Features"]},
        {"num": "02", "title": "GEO", "cn": "生成式引擎优化", "featured": True,
         "desc": "Generative engine citations",
         "subs": ["AI Citation Audit", "Answer Content", "Structured Data", "LLM Monitoring"]},
        {"num": "03", "title": "Campaign", "cn": "整合营销",
         "desc": "Integrated activation",
         "subs": ["Big Idea", "Paid Media", "Social & KOL", "Performance"]},
        {"num": "04", "title": "Strategy", "cn": "品牌战略",
         "desc": "Growth & positioning",
         "subs": ["Market Research", "Positioning", "GTM Strategy", "Growth Roadmap"]},
    ]
    gap, left, top, w, h = 0.28, CHROME_PAD, 1.55, (SLIDE_W - 2 * CHROME_PAD - 0.28 * 3) / 4, 5.0
    for i, svc in enumerate(services[:4]):
        x = left + i * (w + gap)
        featured = svc.get("featured", False)
        bg = BRAND_DEEP if featured else PAPER_2
        tc = WHITE if featured else INK
        sub_c = BRAND_GLOW if featured else INK_3
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(h))
        fill(card, bg)
        card.line.color.rgb = BRAND_DEEP if featured else LINE
        card.line.width = Pt(1)
        label(slide, x + 0.15, top + 0.15, w - 0.3, f"{svc.get('num', '')} · {svc.get('title', '')}", color=sub_c, size=9, mono=True)
        cn = svc.get("cn", "")
        label(slide, x + 0.15, top + 0.42, w - 0.3, cn, color=sub_c, size=10, bold=False)
        label(slide, x + 0.15, top + 0.72, w - 0.3, svc.get("desc", ""), color=tc, size=11, bold=True)
        bullet_box(slide, x + 0.12, top + 1.15, w - 0.24, h - 1.35, svc.get("subs", [])[:4],
                   shade=bg, size=9, text_color=WHITE if featured else INK_3)
    add_slide_footer(slide, sheet_num=sheet_num, total=total)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def proposal_tiers_slide(prs, title="Investment options",
                         tiers=None, logo_path=None, sheet_num=19, total=DEFAULT_TOTAL):
    """Slide 19 — 3 pricing tiers (middle featured)."""
    slide = title_only_slide(prs)
    _slide_bg(slide, PAPER)
    add_slide_header(slide, 19, "PROPOSAL", section_title=title)
    tiers = tiers or [
        {"name": "SEO Audit", "price": "$28K+", "badge": "ONE-TIME",
         "items": ["Technical audit", "Keyword analysis", "Content recs", "GEO snapshot"]},
        {"name": "Full Campaign", "price": "$180K+", "badge": "FEATURED", "featured": True,
         "items": ["Brand strategy", "SEO rebuild", "GEO deploy", "20-week activation", "90-day support"]},
        {"name": "GEO Retainer", "price": "$32K/mo", "badge": "ANNUAL",
         "items": ["Ongoing SEO/GEO", "AI citation monitoring", "Live dashboard", "24h response"]},
    ]
    gap, left, top, w, h = 0.35, CHROME_PAD, 1.55, (SLIDE_W - 2 * CHROME_PAD - 0.7) / 3, 5.0
    for i, tier in enumerate(tiers[:3]):
        x = left + i * (w + gap)
        featured = tier.get("featured", False)
        bg = BRAND_DEEP if featured else PAPER_2
        tc = WHITE if featured else INK
        sub = BRAND_GLOW if featured else INK_3
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(h))
        fill(card, bg)
        card.line.color.rgb = BRAND_DEEP if featured else LINE
        card.line.width = Pt(1)
        if tier.get("badge"):
            label(slide, x + 0.15, top + 0.15, w - 0.3, tier["badge"], color=sub, size=8, mono=True)
        label(slide, x + 0.15, top + 0.42, w - 0.3, tier.get("name", ""), color=tc, size=14, bold=True)
        vtb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top + 0.85), Inches(w - 0.3), Inches(0.65))
        vtf = vtb.text_frame
        no_autofit(vtf)
        vtf.text = tier.get("price", "")
        _style(vtf, tc, 32, bold=True)
        bullet_box(slide, x + 0.12, top + 1.65, w - 0.24, h - 1.85, tier.get("items", [])[:5],
                   shade=bg, size=10, text_color=WHITE if featured else INK_3)
    add_slide_footer(slide, sheet_num=sheet_num, total=total)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def seo_geo_data_slide(prs, metrics=None, message=None, logo_path=None,
                       sheet_num=9, total=DEFAULT_TOTAL):
    """Slide 09 — SEO/GEO industry KPIs with search-specific defaults."""
    default_metrics = [
        {"value": "↓34.5%", "label": "AI Overview CTR", "note": "organic impact", "badge": "Google"},
        {"value": "3.7B", "label": "ChatGPT queries", "note": "weekly volume"},
        {"value": "58%", "label": "Zero-click", "note": "search sessions"},
        {"value": "4.0×", "label": "Precision lift", "note": "GEO vs baseline", "badge": "↑"},
    ]
    return metrics_3up_slide(
        prs, "Search landscape by the numbers", metrics or default_metrics,
        message=message or "AI answers are reshaping how brands earn visibility.",
        logo_path=logo_path, section_num=9, sheet_num=sheet_num, total=total,
    )
